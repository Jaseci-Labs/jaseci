import tiktoken
import docx2txt
import csv
import pptx
from openai import OpenAI
from pypdf import PdfReader
from magic import from_buffer
from tenacity import retry, wait_random_exponential, stop_after_attempt

tokenizer = tiktoken.get_encoding("cl100k_base")


@retry(wait=wait_random_exponential(min=1, max=20), stop=stop_after_attempt(3))
def get_embeddings(texts: list, oai_client: OpenAI, config: dict = {}):
    return [
        result.embedding
        for result in oai_client.embeddings.create(input=texts, **config).data
    ]


def extraction(buff) -> str:
    mimetype = from_buffer(buff.read(), mime=True)
    buff.seek(0)

    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(buff)
        return " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text buff
        return buff.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        return docx2txt.process(buff)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in buff)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
        return extracted_text
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(buff)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
        return extracted_text
    else:
        # Unsupported file type
        raise ValueError("Unsupported file type: {}".format(mimetype))


def extract_text_from_file(file) -> str:
    with open(file, "rb") as buff:
        return " ".join(extraction(buff).split())


def generate_chunks(doc: dict, config: dict) -> list:
    # Check if the document text is empty or whitespace
    text: str = doc.get("text")
    if not text or text.isspace():
        return []

    tokens = tokenizer.encode(text, disallowed_special=())

    # Initialize an empty list of chunks
    chunks = []

    # Use the provided chunk token size or the default one
    chunk_size = config.get("chunk_size") or 200
    max_num_chunks = config.get("max_num_chunks") or 10000
    min_chunk_size_chars = config.get("min_chunk_size_chars") or 350
    min_chunk_length_to_embed = config.get("min_chunk_length_to_embed") or 5

    # Initialize a counter for the number of chunks
    num_chunks = 0

    # Loop until all tokens are consumed
    while tokens and num_chunks < max_num_chunks:
        # Take the first chunk_size tokens as a chunk
        chunk = tokens[:chunk_size]

        # Decode the chunk into text
        chunk_text = tokenizer.decode(chunk)

        # Skip the chunk if it is empty or whitespace
        if not chunk_text or chunk_text.isspace():
            # Remove the tokens corresponding to the chunk text from the remaining tokens
            tokens = tokens[len(chunk) :]
            # Continue to the next iteration of the loop
            continue

        # Find the last period or punctuation mark in the chunk
        last_punctuation = max(
            chunk_text.rfind("."),
            chunk_text.rfind("?"),
            chunk_text.rfind("!"),
            chunk_text.rfind("\n"),
        )

        # If there is a punctuation mark, and the last punctuation index is before MIN_CHUNK_SIZE_CHARS
        if last_punctuation != -1 and last_punctuation > min_chunk_size_chars:
            # Truncate the chunk text at the punctuation mark
            chunk_text = chunk_text[: last_punctuation + 1]

        # Remove any newline characters and strip any leading or trailing whitespace
        chunk_text_to_append = chunk_text.replace("\n", " ").strip()

        if len(chunk_text_to_append) > min_chunk_length_to_embed:
            # Append the cloned doc with chunk text
            chunks.append(clone_doc(doc, chunk_text_to_append))

        # Remove the tokens corresponding to the chunk text from the remaining tokens
        tokens = tokens[len(tokenizer.encode(chunk_text, disallowed_special=())) :]

        # Increment the number of chunks
        num_chunks += 1

    # Handle the remaining tokens
    if tokens:
        remaining_text = tokenizer.decode(tokens).replace("\n", " ").strip()
        if len(remaining_text) > min_chunk_length_to_embed:
            chunks.append(clone_doc(doc, remaining_text))

    return chunks


def clone_doc(doc: dict, text: str) -> dict:
    _doc = doc.copy()
    _doc["text"] = text
    return _doc
